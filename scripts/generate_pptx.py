#!/usr/bin/env python3
"""
0xH (헥세이치) — PPT 자동 생성 스크립트

사용법:
    python scripts/generate_pptx.py slides.json -o output.pptx

slides.json 포맷은 scripts/slide_data_example.py 또는
docs/templates/slide-package-template.md 참고.
"""

import json
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ─── 0xH 브랜드 컬러 ─────────────────────────────────────
COLORS = {
    "bg_dark": RGBColor(0x0F, 0x0F, 0x23),        # #0f0f23 배경
    "bg_card": RGBColor(0x1A, 0x1A, 0x2E),         # #1a1a2e 카드/박스
    "teal": RGBColor(0x0D, 0x94, 0x88),             # #0d9488 Cyber Teal
    "teal_dark": RGBColor(0x0F, 0x76, 0x6E),        # #0f766e Deep Teal
    "text_primary": RGBColor(0xE0, 0xE0, 0xE0),     # #e0e0e0 밝은 회색
    "text_secondary": RGBColor(0x9C, 0xA3, 0xAF),   # #9ca3af 보조 텍스트
    "error": RGBColor(0xEF, 0x44, 0x44),             # #ef4444 에러/문제
    "success": RGBColor(0x22, 0xC5, 0x5E),           # #22c55e 성공/해결
    "warning": RGBColor(0xEA, 0xB3, 0x08),           # #eab308 주의
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "code_bg": RGBColor(0x14, 0x14, 0x28),           # 코드 블록 배경
}

# 폰트 설정
FONT_TITLE = "Pretendard"       # 없으면 맑은 고딕 사용
FONT_BODY = "Pretendard"
FONT_CODE = "JetBrains Mono"    # 없으면 Consolas 사용
FONT_FALLBACK_TITLE = "맑은 고딕"
FONT_FALLBACK_BODY = "맑은 고딕"
FONT_FALLBACK_CODE = "Consolas"


def set_slide_bg(slide, color):
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    """색이 채워진 둥근 사각형 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # 테두리 없음
    # 둥근 정도 설정
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_name=None, color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name or FONT_BODY
    p.font.color.rgb = color or COLORS["text_primary"]
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=18,
                       font_name=None, color=None, line_spacing=1.5):
    """여러 줄 텍스트 박스 추가. lines는 문자열 리스트 또는 (text, color, bold) 튜플 리스트."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, tuple):
            text, line_color, bold = line
        else:
            text, line_color, bold = line, color or COLORS["text_primary"], False

        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name or FONT_BODY
        p.font.color.rgb = line_color
        p.font.bold = bold
        p.space_after = Pt(font_size * (line_spacing - 1))

    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=14):
    """코드 블록 (다크 배경 박스 + 모노 폰트)."""
    # 배경 박스
    bg_shape = add_shape(slide, left, top, width, height, COLORS["code_bg"])

    # 코드 텍스트
    txBox = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(0.2),
        width - Inches(0.6), height - Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = FONT_CODE
        p.font.color.rgb = COLORS["text_primary"]
        p.space_after = Pt(2)

    return bg_shape


def add_table(slide, left, top, width, rows_data, col_widths=None, font_size=14):
    """테이블 추가. rows_data: [["header1", "header2"], ["val1", "val2"], ...]"""
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 2
    height = Inches(0.4 * rows + 0.1)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 열 너비 설정
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            # 셀 스타일
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT_BODY
                if row_idx == 0:  # 헤더
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLORS["white"]
                else:
                    paragraph.font.color.rgb = COLORS["text_primary"]

            # 셀 배경
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = COLORS["teal_dark"]
            else:
                cell_fill.fore_color.rgb = COLORS["bg_card"] if row_idx % 2 == 1 else COLORS["bg_dark"]

    return table_shape


def set_speaker_notes(slide, notes_text):
    """슬라이드 스피커 노트 설정."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def build_title_slide(prs, data):
    """타이틀 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, COLORS["bg_dark"])

    # Teal 악센트 라인 (상단)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["teal"]
    accent.line.fill.background()

    # 제목
    add_text_box(
        slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
        data.get("title", "제목"),
        font_size=40, color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # 부제목
    subtitle = data.get("subtitle", "")
    if subtitle:
        add_text_box(
            slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
            subtitle,
            font_size=22, color=COLORS["teal"],
            alignment=PP_ALIGN.CENTER
        )

    # 브랜드
    add_text_box(
        slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
        "0xH — 헥세이치",
        font_size=18, color=COLORS["text_secondary"],
        alignment=PP_ALIGN.CENTER
    )

    # 스피커 노트
    if "notes_ko" in data:
        notes = f"[KO]\n{data['notes_ko']}"
        if "notes_en" in data:
            notes += f"\n\n[EN]\n{data['notes_en']}"
        set_speaker_notes(slide, notes)


def build_toc_slide(prs, data):
    """목차 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_dark"])

    # 섹션 제목
    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
        "오늘 다룰 내용",
        font_size=32, color=COLORS["white"], bold=True
    )

    # Teal 밑줄
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.2), Inches(2), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["teal"]
    accent.line.fill.background()

    # 목차 항목들
    items = data.get("items", [])
    for i, item in enumerate(items):
        y = 1.8 + i * 1.0

        # 번호 (Teal 원)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), Inches(y), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["teal"]
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = COLORS["white"]
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 항목 텍스트
        add_text_box(
            slide, Inches(1.8), Inches(y + 0.05), Inches(9), Inches(0.5),
            item,
            font_size=22, color=COLORS["text_primary"]
        )

    # 스피커 노트
    if "notes_ko" in data:
        notes = f"[KO]\n{data['notes_ko']}"
        if "notes_en" in data:
            notes += f"\n\n[EN]\n{data['notes_en']}"
        set_speaker_notes(slide, notes)


def build_content_slide(prs, data):
    """일반 콘텐츠 슬라이드 (제목 + 본문 텍스트)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_dark"])

    # 섹션 번호 + 제목
    section_label = data.get("section", "")
    title = data.get("title", "")

    if section_label:
        add_text_box(
            slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.5),
            section_label,
            font_size=14, color=COLORS["teal"], bold=True
        )

    add_text_box(
        slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
        title,
        font_size=30, color=COLORS["white"], bold=True
    )

    # Teal 밑줄
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["teal"]
    accent.line.fill.background()

    # 본문 내용
    y_offset = 1.8
    content = data.get("content", [])

    for block in content:
        block_type = block.get("type", "text")

        if block_type == "text":
            lines = block.get("lines", [])
            processed_lines = []
            for line in lines:
                if isinstance(line, dict):
                    processed_lines.append((
                        line["text"],
                        _parse_color(line.get("color", "text_primary")),
                        line.get("bold", False)
                    ))
                else:
                    processed_lines.append(line)

            box = add_multiline_text(
                slide, Inches(0.8), Inches(y_offset), Inches(11), Inches(3),
                processed_lines, font_size=20
            )
            y_offset += len(processed_lines) * 0.45

        elif block_type == "code":
            code = block.get("code", "")
            code_height = min(len(code.split("\n")) * 0.35 + 0.4, 3.5)
            add_code_block(
                slide, Inches(0.8), Inches(y_offset),
                Inches(11.4), Inches(code_height),
                code, font_size=block.get("font_size", 16)
            )
            y_offset += code_height + 0.2

        elif block_type == "table":
            rows = block.get("rows", [])
            col_widths = block.get("col_widths")
            if col_widths:
                col_widths = [Inches(w) for w in col_widths]
            add_table(
                slide, Inches(0.8), Inches(y_offset),
                Inches(11.4), rows,
                col_widths=col_widths,
                font_size=block.get("font_size", 16)
            )
            y_offset += len(rows) * 0.4 + 0.3

        elif block_type == "badge":
            # 결과 뱃지 (성공/실패/부분성공)
            status = block.get("status", "")
            badge_colors = {
                "success": COLORS["success"],
                "error": COLORS["error"],
                "warning": COLORS["warning"],
            }
            badge_color = badge_colors.get(block.get("color", ""), COLORS["text_secondary"])
            badge_shape = add_shape(
                slide, Inches(0.8), Inches(y_offset),
                Inches(11.4), Inches(0.6), COLORS["bg_card"]
            )
            add_text_box(
                slide, Inches(1.1), Inches(y_offset + 0.1),
                Inches(10.8), Inches(0.5),
                status,
                font_size=20, color=badge_color, bold=True,
                alignment=PP_ALIGN.LEFT
            )
            y_offset += 0.8

    # 스피커 노트
    if "notes_ko" in data:
        notes = f"[KO]\n{data['notes_ko']}"
        if "notes_en" in data:
            notes += f"\n\n[EN]\n{data['notes_en']}"
        set_speaker_notes(slide, notes)


def build_screen_recording_slide(prs, data):
    """화면녹화 전환 슬라이드 (녹화 구간 안내)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_dark"])

    # 중앙 안내 박스
    box = add_shape(
        slide, Inches(3), Inches(2.5), Inches(7), Inches(2.5),
        COLORS["bg_card"]
    )

    add_text_box(
        slide, Inches(3.5), Inches(2.8), Inches(6), Inches(0.5),
        "화면 녹화",
        font_size=14, color=COLORS["teal"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(3.5), Inches(3.3), Inches(6), Inches(1.2),
        data.get("description", "[화면녹화 내용 설명]"),
        font_size=22, color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )

    # 스피커 노트
    if "notes_ko" in data:
        notes = f"[KO]\n{data['notes_ko']}"
        if "notes_en" in data:
            notes += f"\n\n[EN]\n{data['notes_en']}"
        set_speaker_notes(slide, notes)


def build_outro_slide(prs, data):
    """아웃트로 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_dark"])

    # Teal 악센트 라인 (하단)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.44), prs.slide_width, Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["teal"]
    accent.line.fill.background()

    # 블로그 링크
    blog_url = data.get("blog_url", "techblips.com")
    add_text_box(
        slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
        f"📝 블로그: {blog_url}",
        font_size=22, color=COLORS["teal"],
        alignment=PP_ALIGN.CENTER
    )

    # 다음 영상
    next_topic = data.get("next_topic", "")
    if next_topic:
        add_text_box(
            slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
            f"다음 영상: {next_topic}",
            font_size=22, color=COLORS["text_primary"],
            alignment=PP_ALIGN.CENTER
        )

    # 구독 CTA
    cta_box = add_shape(
        slide, Inches(4.5), Inches(4), Inches(4), Inches(0.8),
        COLORS["teal"]
    )
    add_text_box(
        slide, Inches(4.5), Inches(4.1), Inches(4), Inches(0.6),
        "구독 & 좋아요",
        font_size=24, color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # 브랜드
    add_text_box(
        slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
        "0xH — 헥세이치",
        font_size=20, color=COLORS["text_secondary"],
        alignment=PP_ALIGN.CENTER
    )

    # 스피커 노트
    if "notes_ko" in data:
        notes = f"[KO]\n{data['notes_ko']}"
        if "notes_en" in data:
            notes += f"\n\n[EN]\n{data['notes_en']}"
        set_speaker_notes(slide, notes)


def _parse_color(color_name):
    """컬러 이름을 RGBColor로 변환."""
    return COLORS.get(color_name, COLORS["text_primary"])


# ─── 슬라이드 타입 매핑 ──────────────────────────────────
SLIDE_BUILDERS = {
    "title": build_title_slide,
    "toc": build_toc_slide,
    "content": build_content_slide,
    "screen_recording": build_screen_recording_slide,
    "outro": build_outro_slide,
}


def generate_pptx(slides_data, output_path):
    """JSON 데이터로 PPTX 생성."""
    prs = Presentation()
    # 16:9 비율
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data.get("slides", []):
        slide_type = slide_data.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type, build_content_slide)
        builder(prs, slide_data)

    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    print(f"   슬라이드 수: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(description="0xH PPT 자동 생성")
    parser.add_argument("input", help="슬라이드 데이터 JSON 파일")
    parser.add_argument("-o", "--output", default="output.pptx", help="출력 PPTX 파일명")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        data = json.load(f)

    generate_pptx(data, args.output)


if __name__ == "__main__":
    main()
